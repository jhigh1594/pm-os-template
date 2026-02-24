#!/usr/bin/env python3
"""
merge-title-content.py

Combines a title slide from the corporate template with html2pptx-generated content slides.

Usage:
    python merge-title-content.py <title_pptx> <content_pptx> <output_pptx>

Example:
    python merge-title-content.py /tmp/planview-deck/with-title.pptx /tmp/planview-deck/content.pptx /tmp/planview-deck/final.pptx
"""

import sys
import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt


def copy_slide(source_prs, source_idx, target_prs):
    """Copy a slide from source presentation to target presentation."""
    source_slide = source_prs.slides[source_idx]

    # Get the blank layout from target
    blank_layout = target_prs.slide_layouts[6]  # Usually blank layout

    # Create new slide
    new_slide = target_prs.slides.add_slide(blank_layout)

    # Copy shapes
    for shape in source_slide.shapes:
        if shape.has_text_frame:
            # Copy text box
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            textbox = new_slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame

            # Copy paragraphs
            for i, para in enumerate(shape.text_frame.paragraphs):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = para.text
                p.alignment = para.alignment

                # Copy font properties
                if para.runs:
                    for j, run in enumerate(para.runs):
                        if j == 0:
                            r = p.runs[0] if p.runs else p.add_run()
                        else:
                            r = p.add_run()
                        r.text = run.text
                        if run.font.size:
                            r.font.size = run.font.size
                        if run.font.bold is not None:
                            r.font.bold = run.font.bold
                        if run.font.italic is not None:
                            r.font.italic = run.font.italic
                        if run.font.color.rgb:
                            r.font.color.rgb = run.font.color.rgb

        elif hasattr(shape, 'image'):
            # Copy image
            try:
                image = shape.image
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Save image to temp and re-add
                image_bytes = image.blob
                ext = image.ext

                import tempfile
                with tempfile.NamedTemporaryFile(suffix=f'.{ext}', delete=False) as f:
                    f.write(image_bytes)
                    temp_path = f.name

                new_slide.shapes.add_picture(temp_path, left, top, width, height)
                os.unlink(temp_path)
            except Exception as e:
                print(f"Warning: Could not copy image: {e}")

    return new_slide


def merge_presentations(title_path, content_path, output_path):
    """Merge title slide with content slides."""
    print(f"Loading title presentation: {title_path}")
    title_prs = Presentation(title_path)

    print(f"Loading content presentation: {content_path}")
    content_prs = Presentation(content_path)

    # Create new presentation with same dimensions
    output_prs = Presentation()
    output_prs.slide_width = content_prs.slide_width
    output_prs.slide_height = content_prs.slide_height

    # Copy title slide (slide 0) from title presentation
    print("Copying title slide...")
    copy_slide(title_prs, 0, output_prs)

    # Copy all slides from content presentation
    print(f"Copying {len(content_prs.slides)} content slides...")
    for i in range(len(content_prs.slides)):
        copy_slide(content_prs, i, output_prs)

    # Save output
    print(f"Saving to: {output_path}")
    output_prs.save(output_path)

    print(f"Done! Created presentation with {len(output_prs.slides)} slides")


def main():
    if len(sys.argv) != 4:
        print(__doc__)
        sys.exit(1)

    title_path = sys.argv[1]
    content_path = sys.argv[2]
    output_path = sys.argv[3]

    if not os.path.exists(title_path):
        print(f"Error: Title file not found: {title_path}")
        sys.exit(1)

    if not os.path.exists(content_path):
        print(f"Error: Content file not found: {content_path}")
        sys.exit(1)

    merge_presentations(title_path, content_path, output_path)


if __name__ == "__main__":
    main()
