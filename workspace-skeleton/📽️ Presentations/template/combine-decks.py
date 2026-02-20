#!/usr/bin/env python3
"""
Combine two PPTX files by copying slides from content to title template.
Preserves master slides from title template for consistent branding.

Usage: python combine-decks.py title.pptx content.pptx output.pptx
"""

import sys
import os
from io import BytesIO
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def copy_slide(source_prs, source_idx, target_prs):
    """
    Copy a slide from source to target presentation.
    This copies the slide layout and all shapes.
    """
    source_slide = source_prs.slides[source_idx]

    # Get the slide layout from the source slide
    # Use blank layout in target as base
    blank_layout = target_prs.slide_layouts[6]  # Usually blank
    new_slide = target_prs.slides.add_slide(blank_layout)

    # Copy each shape
    for shape in source_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Copy image
            try:
                image = shape.image
                image_bytes = BytesIO(image.blob)  # Wrap in BytesIO
                new_slide.shapes.add_picture(
                    image_bytes,
                    shape.left, shape.top,
                    shape.width, shape.height
                )
            except Exception as e:
                print(f"  Warning: Could not copy image: {e}")

        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.has_text_frame:
            # Copy text box
            try:
                txBox = new_slide.shapes.add_textbox(
                    shape.left, shape.top,
                    shape.width, shape.height
                )
                tf = txBox.text_frame

                # Copy text frame properties
                tf.word_wrap = shape.text_frame.word_wrap

                # Copy paragraphs
                for i, para in enumerate(shape.text_frame.paragraphs):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    # Don't set p.text - we'll copy via runs instead
                    p.alignment = para.alignment
                    p.level = para.level

                    # Clear the default run and copy runs with font properties
                    if para.runs:
                        # Clear default run text first
                        if p.runs:
                            p.runs[0].text = ''

                        for j, run in enumerate(para.runs):
                            if j == 0:
                                r = p.runs[0]
                            else:
                                r = p.add_run()

                            r.text = run.text

                        # Safely copy font properties
                        try:
                            if run.font.size:
                                r.font.size = run.font.size
                        except:
                            pass

                        try:
                            if run.font.bold is not None:
                                r.font.bold = run.font.bold
                        except:
                            pass

                        try:
                            if run.font.italic is not None:
                                r.font.italic = run.font.italic
                        except:
                            pass

                        try:
                            if run.font.name:
                                r.font.name = run.font.name
                        except:
                            pass

                        # Handle font color carefully
                        try:
                            if run.font.color and run.font.color.type is not None:
                                if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                    r.font.color.rgb = run.font.color.rgb
                        except:
                            pass

            except Exception as e:
                print(f"  Warning: Could not copy text box: {e}")

        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Copy auto shape (rectangles, etc.)
            try:
                # Get shape name/type
                shape_name = shape.shape_type.name if hasattr(shape.shape_type, 'name') else 'rect'

                if 'RECTANGLE' in str(shape_name).upper() or shape.shape_type == MSO_SHAPE_TYPE.RECTANGLE:
                    new_shape = new_slide.shapes.add_shape(
                        1,  # Rectangle
                        shape.left, shape.top,
                        shape.width, shape.height
                    )
                    # Copy fill
                    if shape.fill.type is not None:
                        new_shape.fill.solid()
                        if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color.rgb:
                            new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
            except Exception as e:
                print(f"  Warning: Could not copy shape: {e}")

        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # Copy table
            try:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)
                new_table = new_slide.shapes.add_table(
                    rows, cols,
                    shape.left, shape.top,
                    shape.width, shape.height
                ).table

                # Copy cell contents
                for r in range(rows):
                    for c in range(cols):
                        cell = table.cell(r, c)
                        new_cell = new_table.cell(r, c)
                        new_cell.text = cell.text
            except Exception as e:
                print(f"  Warning: Could not copy table: {e}")

def combine_presentations(title_path, content_path, output_path):
    """
    Combine title slide with content slides.
    Uses title presentation as base to preserve master slides and branding.
    """
    # Load both presentations
    print(f"Loading title presentation: {title_path}")
    title_prs = Presentation(title_path)

    print(f"Loading content presentation: {content_path}")
    content_prs = Presentation(content_path)

    print(f"\nTitle presentation has {len(title_prs.slides)} slide(s)")
    print(f"Content presentation has {len(content_prs.slides)} slide(s)")

    # Copy content slides to title presentation
    for i in range(len(content_prs.slides)):
        print(f"  Copying content slide {i + 1}...")
        copy_slide(content_prs, i, title_prs)

    # Save combined presentation
    title_prs.save(output_path)
    print(f"\nCombined presentation saved: {output_path}")
    print(f"  - {1} title slide")
    print(f"  - {len(content_prs.slides)} content slides")
    print(f"  - {len(title_prs.slides)} total slides")

if __name__ == "__main__":
    if len(sys.argv) != 4:
        print("Usage: python combine-decks.py title.pptx content.pptx output.pptx")
        sys.exit(1)

    title_path = sys.argv[1]
    content_path = sys.argv[2]
    output_path = sys.argv[3]

    if not os.path.exists(title_path):
        print(f"Error: Title file not found: {title_path}")
        sys.exit(1)

    if not os.path.exists(content_path):
        print(f"Error: Content file not found: {content_path}")
        sys.exit(1)

    combine_presentations(title_path, content_path, output_path)
